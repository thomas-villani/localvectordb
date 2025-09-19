# Copyright (c) 2023-2025 Tom Villani, Ph.D.
#
# This work is licensed under the Creative Commons Attribution-NonCommercial 4.0 International License.
# You may not use this file for commercial purposes without explicit permission.
#
# For more information, please visit: https://creativecommons.org/licenses/by-nc/4.0/
#
# Contact: thomas.villani@gmail.com
#
# src/localvectordb/extractors/office_extractors.py
# src/localvectordb/extractors/office_extractors.py
"""
Microsoft Office document extractors.
"""

import io
import logging
from typing import List, Optional

from localvectordb import MetadataField
from localvectordb.core import MetadataFieldType
from localvectordb.extractors import BaseExtractor, ExtractionResult

logger = logging.getLogger(__name__)


class DocxExtractor(BaseExtractor):
    """
    Extractor for Microsoft Word documents (.docx).
    """

    @property
    def supported_extensions(self) -> List[str]:
        return ['.docx']

    @property
    def supported_mimetypes(self) -> List[str]:
        return [
            'application/vnd.openxmlformats-officedocument.wordprocessingml.document'
        ]

    @property
    def required_packages(self) -> List[str]:
        return ['python-docx']

    @property
    def priority(self) -> int:
        return 15

    @property
    def metadata_schema(self) -> dict[str, MetadataField]:
        return {
            "filename": MetadataField(type=MetadataFieldType.TEXT, indexed=True, required=False),
            "paragraph_count": MetadataField(type=MetadataFieldType.INTEGER, indexed=False, required=False),
            "table_count": MetadataField(type=MetadataFieldType.INTEGER, indexed=False, required=False),
            "file_size_bytes": MetadataField(type=MetadataFieldType.INTEGER, indexed=False, required=False),
            "character_count": MetadataField(type=MetadataFieldType.INTEGER, indexed=False, required=False)
        }

    def _check_availability(self) -> bool:
        try:
            import docx
            return True
        except ImportError:
            return False

    def _extract_text_impl(
            self, file_content: bytes, filename: str, mimetype: Optional[str], **kwargs
    ) -> ExtractionResult:
        """Extract text from DOCX files."""
        try:
            from docx import Document

            with io.BytesIO(file_content) as file_buffer:
                doc = Document(file_buffer)

                text_parts = []
                paragraph_count = 0

                # Extract text from paragraphs
                for paragraph in doc.paragraphs:
                    paragraph_text = paragraph.text.strip()
                    if paragraph_text:
                        text_parts.append(paragraph_text)
                        paragraph_count += 1

                # Extract text from tables
                table_count = 0
                for table in doc.tables:
                    table_count += 1
                    table_text = []
                    for row in table.rows:
                        row_text = []
                        for cell in row.cells:
                            cell_text = cell.text.strip()
                            if cell_text:
                                row_text.append(cell_text)
                        if row_text:
                            table_text.append('\t'.join(row_text))

                    if table_text:
                        text_parts.append(f"\n[Table {table_count}]\n" + '\n'.join(table_text))

                if not text_parts:
                    return ExtractionResult(
                        text="",
                        success=False,
                        method='DocxExtractor',
                        error="No text content found in Word document"
                    )

                full_text = '\n\n'.join(text_parts)

                metadata = {
                    'filename': filename,
                    'paragraph_count': paragraph_count,
                    'table_count': table_count,
                    'file_size_bytes': len(file_content),
                    'character_count': len(full_text)
                }

                return ExtractionResult(
                    text=full_text,
                    success=True,
                    method='DocxExtractor',
                    metadata=metadata
                )

        except Exception as e:
            return ExtractionResult(
                text="",
                success=False,
                method='DocxExtractor',
                error=f"DOCX extraction failed: {str(e)}"
            )


class PptxExtractor(BaseExtractor):
    """
    Extractor for Microsoft PowerPoint presentations (.pptx).
    """

    @property
    def supported_extensions(self) -> List[str]:
        return ['.pptx']

    @property
    def supported_mimetypes(self) -> List[str]:
        return [
            'application/vnd.openxmlformats-officedocument.presentationml.presentation'
        ]

    @property
    def required_packages(self) -> List[str]:
        return ['python-pptx']

    @property
    def priority(self) -> int:
        return 15

    @property
    def metadata_schema(self) -> dict[str, MetadataField]:
        return {
            "filename": MetadataField(type=MetadataFieldType.TEXT, indexed=True, required=False),
            "slide_count": MetadataField(type=MetadataFieldType.INTEGER, indexed=False, required=False),
            "slides_with_text": MetadataField(type=MetadataFieldType.INTEGER, indexed=False, required=False),
            "total_shapes": MetadataField(type=MetadataFieldType.INTEGER, indexed=False, required=False),
            "file_size_bytes": MetadataField(type=MetadataFieldType.INTEGER, indexed=False, required=False),
            "character_count": MetadataField(type=MetadataFieldType.INTEGER, indexed=False, required=False),
        }

    def _check_availability(self) -> bool:
        try:
            import pptx
            return True
        except ImportError:
            return False

    def _extract_text_impl(
            self, file_content: bytes, filename: str, mimetype: Optional[str], **kwargs
    ) -> ExtractionResult:
        """Extract text from PPTX files."""
        try:
            from pptx import Presentation

            with io.BytesIO(file_content) as file_buffer:
                prs = Presentation(file_buffer)

                slide_texts = []
                total_shapes = 0

                for slide_num, slide in enumerate(prs.slides, 1):
                    slide_content = []
                    shapes_with_text = 0

                    for shape in slide.shapes:
                        total_shapes += 1

                        # Extract text from text shapes
                        if hasattr(shape, "text") and shape.text.strip():
                            slide_content.append(shape.text.strip())
                            shapes_with_text += 1

                        # Extract text from tables in slides
                        elif hasattr(shape, "table"):
                            table_text = []
                            for row in shape.table.rows:
                                row_text = []
                                for cell in row.cells:
                                    cell_text = cell.text.strip()
                                    if cell_text:
                                        row_text.append(cell_text)
                                if row_text:
                                    table_text.append('\t'.join(row_text))

                            if table_text:
                                slide_content.append("[Table]\n" + '\n'.join(table_text))
                                shapes_with_text += 1

                    if slide_content:
                        slide_text = f"[Slide {slide_num}]\n" + '\n'.join(slide_content)
                        slide_texts.append(slide_text)

                if not slide_texts:
                    return ExtractionResult(
                        text="",
                        success=False,
                        method='PptxExtractor',
                        error="No text content found in PowerPoint presentation"
                    )

                full_text = '\n\n---\n\n'.join(slide_texts)

                metadata = {
                    'filename': filename,
                    'slide_count': len(prs.slides),
                    'slides_with_text': len(slide_texts),
                    'total_shapes': total_shapes,
                    'file_size_bytes': len(file_content),
                    'character_count': len(full_text)
                }

                return ExtractionResult(
                    text=full_text,
                    success=True,
                    method='PptxExtractor',
                    metadata=metadata
                )

        except Exception as e:
            return ExtractionResult(
                text="",
                success=False,
                method='PptxExtractor',
                error=f"PPTX extraction failed: {str(e)}"
            )


class XlsxExtractor(BaseExtractor):
    """
    Extractor for Microsoft Excel spreadsheets (.xlsx).
    """

    @property
    def supported_extensions(self) -> List[str]:
        return ['.xlsx']

    @property
    def supported_mimetypes(self) -> List[str]:
        return [
            'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet'
        ]

    @property
    def required_packages(self) -> List[str]:
        return ['openpyxl']

    @property
    def priority(self) -> int:
        return 15

    @property
    def metadata_schema(self) -> dict[str, MetadataField]:
        return {
            "filename": MetadataField(type=MetadataFieldType.TEXT, indexed=True, required=False),
            "sheet_count": MetadataField(type=MetadataFieldType.INTEGER, indexed=False, required=False),
            "sheets_with_data": MetadataField(type=MetadataFieldType.INTEGER, indexed=False, required=False),
            "total_cells": MetadataField(type=MetadataFieldType.INTEGER, indexed=False, required=False),
            "non_empty_cells": MetadataField(type=MetadataFieldType.INTEGER, indexed=False, required=False),
            "file_size_bytes": MetadataField(type=MetadataFieldType.INTEGER, indexed=False, required=False),
            "character_count": MetadataField(type=MetadataFieldType.INTEGER, indexed=False, required=False),
        }

    def _check_availability(self) -> bool:
        try:
            import openpyxl
            return True
        except ImportError:
            return False

    def _extract_text_impl(
            self, file_content: bytes, filename: str, mimetype: Optional[str], **kwargs
    ) -> ExtractionResult:
        """Extract text from XLSX files."""
        try:
            from openpyxl import load_workbook

            with io.BytesIO(file_content) as file_buffer:
                wb = load_workbook(file_buffer, data_only=True)

                sheet_texts = []
                total_cells = 0
                non_empty_cells = 0

                for sheet_name in wb.sheetnames:
                    sheet = wb[sheet_name]
                    sheet_content = [f"[Sheet: {sheet_name}]"]

                    # Get the dimensions of the sheet
                    max_row = sheet.max_row
                    max_col = sheet.max_column

                    if max_row == 1 and max_col == 1:
                        # Empty sheet
                        continue

                    # Extract data row by row
                    for row_num, row in enumerate(sheet.iter_rows(values_only=True), 1):
                        row_data = []

                        for cell_value in row:
                            total_cells += 1
                            if cell_value is not None:
                                # Convert to string and clean up
                                cell_str = str(cell_value).strip()
                                if cell_str:
                                    row_data.append(cell_str)
                                    non_empty_cells += 1
                                else:
                                    row_data.append("")
                            else:
                                row_data.append("")

                        # Only add rows that have some content
                        if any(cell.strip() for cell in row_data):
                            sheet_content.append('\t'.join(row_data))

                    # Only add sheets that have content beyond the header
                    if len(sheet_content) > 1:
                        sheet_texts.append('\n'.join(sheet_content))

                if not sheet_texts:
                    return ExtractionResult(
                        text="",
                        success=False,
                        method='XlsxExtractor',
                        error="No text content found in Excel spreadsheet"
                    )

                full_text = '\n\n'.join(sheet_texts)

                metadata = {
                    'filename': filename,
                    'sheet_count': len(wb.sheetnames),
                    'sheets_with_data': len(sheet_texts),
                    'total_cells': total_cells,
                    'non_empty_cells': non_empty_cells,
                    'file_size_bytes': len(file_content),
                    'character_count': len(full_text)
                }

                return ExtractionResult(
                    text=full_text,
                    success=True,
                    method='XlsxExtractor',
                    metadata=metadata
                )

        except Exception as e:
            return ExtractionResult(
                text="",
                success=False,
                method='XlsxExtractor',
                error=f"XLSX extraction failed: {str(e)}"
            )
